"""PPTX parser — extracts structured sections and text from PowerPoint presentations.

Requires the ``pptx`` extra::

    pip install smartchunk[pptx] or pip install python-pptx
"""

from __future__ import annotations

from pathlib import Path

from smartchunk.models import DocumentSection
from smartchunk.parsers.base import BaseParser


class PptxParser(BaseParser):
    """Parses Microsoft PowerPoint (.pptx) documents."""

    supported_extensions = [".pptx"]

    def parse(self, filepath: str | Path) -> list[DocumentSection]:
        filepath = Path(filepath)

        try:
            import pptx  # type: ignore[import-untyped]
        except ImportError as exc:
            raise ImportError(
                "PPTX parsing requires python-pptx. "
                "Install with: pip install python-pptx"
            ) from exc

        prs = pptx.Presentation(str(filepath))
        sections: list[DocumentSection] = []

        for idx, slide in enumerate(prs.slides):
            slide_text_parts: list[str] = []

            # Extract slide title if available
            slide_title = ""
            if slide.shapes.title:
                slide_title = slide.shapes.title.text.strip()

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text and text != slide_title:
                            slide_text_parts.append(text)

            slide_text = "\n".join(slide_text_parts)
            heading = slide_title if slide_title else f"Slide {idx + 1}"

            sections.append(
                DocumentSection(
                    text=slide_text,
                    heading=heading,
                    level=1,
                    parent_headings=[],
                    metadata={"source": filepath.name, "slide_index": idx},
                )
            )

        return sections
